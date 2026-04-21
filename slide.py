#!/usr/bin/env python3
"""
Sprint 0 — AI Usage Report Slide Generator

Reads the best-practices markdown report, uses LLM (gpt-5.4-mini) to
condense each section into presentation-ready bullet points, then builds
a branded 16:9 .pptx with CMC Global theme + chart images.

Usage:
    python slide.py                          # default output: sprint_0_report.pptx
    python slide.py -o my_deck.pptx          # custom output path
    python slide.py --model gpt-5.4-mini     # override model
"""

from __future__ import annotations

import argparse
import os
import re
import textwrap
from pathlib import Path
from typing import Any

import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from dotenv import load_dotenv
    load_dotenv()
except Exception:
    pass

# ── Paths ───────────────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent
REPORT_DIR = BASE_DIR / "report_sprint_0_v1_200426"
CHARTS_DIR = REPORT_DIR / "charts_output"
LOGO_PATH = BASE_DIR / "logo.png"


def _find_report_md() -> Path:
    """Search for the best practices markdown report, preferring the report dir."""
    candidates = [
        REPORT_DIR / "best_practices_report_sprint_0_v1.md",
        BASE_DIR / "log" / "best_practices_report_sprint_0_v1.md",
        BASE_DIR / "log" / "best_practices_report_sprint_0.md",
    ]
    for c in candidates:
        if c.exists():
            return c
    # Fallback: glob for any best_practices*.md
    for p in sorted(BASE_DIR.rglob("best_practices*.md"), key=lambda x: x.stat().st_mtime, reverse=True):
        return p
    raise FileNotFoundError("No best practices report (.md) found in project")

# ── Brand colours (CMC Global) ──────────────────────────────────────────────
CMC_BLUE     = RGBColor(0x1A, 0x9F, 0xD9)   # primary accent
CMC_DARK     = RGBColor(0x0D, 0x2C, 0x54)   # heading / dark navy
CMC_LIGHT_BG = RGBColor(0xF2, 0xF8, 0xFC)   # soft background
CMC_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CMC_GRAY     = RGBColor(0x5A, 0x5A, 0x5A)   # body text
CMC_ACCENT2  = RGBColor(0xE8, 0x6C, 0x00)   # orange accent for highlights

# ── OpenAI / Azure helper (reused from generate_report.py) ──────────────────
OPENAI_BASE_URL = os.getenv("OPENAI_BASE_URL")
OPENAI_API_KEY  = os.getenv("OPENAI_API_KEY")


def _call_llm(model: str, prompt: str, timeout: int = 120) -> str:
    base = (OPENAI_BASE_URL or "").rstrip("/")
    is_azure = "cognitiveservices.azure.com" in base or "openai.azure.com" in base

    if is_azure:
        host = re.sub(r"/openai(/v1)?$", "", base).rstrip("/")
        url = f"{host}/openai/deployments/{model}/chat/completions?api-version=2024-12-01-preview"
        headers = {"api-key": OPENAI_API_KEY or ""}
    else:
        url = f"{base}/chat/completions"
        headers = {}
        if OPENAI_API_KEY:
            headers["Authorization"] = f"Bearer {OPENAI_API_KEY}"

    body: dict[str, Any] = {
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
        "max_completion_tokens": 4000,
    }
    if not is_azure:
        body["model"] = model

    r = requests.post(url, headers=headers, json=body, timeout=timeout)
    r.raise_for_status()
    return r.json()["choices"][0]["message"]["content"]


# ── LLM-powered content condenser ──────────────────────────────────────────

SYSTEM_SLIDE = textwrap.dedent("""\
You are a senior project manager preparing a client-facing PowerPoint presentation.
Your audience is NON-TECHNICAL (executives, clients who do not understand coding agents).
Rules:
- Return ONLY bullet points (use "- " prefix), no headings, no markdown formatting.
- Each bullet max 10 words with analysis (always include a numerical value - percentage, number, or specific example). Prefer plain business language.
- Each bullet needs to have a noun phrase, for example: "Documentation: 8 tasks, 50% time saved" or "Tester: 23.5h saved (62% efficiency)".
- Each slide needs to have a summary insights section.
- Focus on the task's performance of each team member using AI tools, not the technical details of the tools themselves.
- Replace jargon: "prompt" → "instruction to AI", "context ordering" → "giving AI clear info".
- No code, no XML, no technical syntax.
- Keep it concise: max {max_bullets} bullets.
""")


def condense(model: str, section_text: str, slide_title: str,
             max_bullets: int = 6) -> list[str]:
    """Send a report section to LLM → get back short bullet points."""
    prompt = (
        SYSTEM_SLIDE.format(max_bullets=max_bullets)
        + f"\n---\nSlide title: {slide_title}\n"
        + f"Source content:\n{section_text}\n---\n"
        + "Return bullet points now."
    )
    raw = _call_llm(model, prompt)
    bullets = []
    for line in raw.splitlines():
        line = line.strip().lstrip("-•*").strip()
        if line:
            bullets.append(line)
    return bullets[:max_bullets]


def generate_conclusion(model: str, full_report: str) -> list[str]:
    """Ask LLM to write a conclusion slide from the full report."""
    prompt = textwrap.dedent(f"""\
    You are a senior project manager writing a conclusion slide for executives.
    Audience: non-technical clients and company leadership.
    Based on the report below, write 5-6 bullet points summarising:
    1. The overall result of the AI pilot (time saved, efficiency %).
    2. The biggest win.
    3. The biggest area for improvement.
    4. What the team should do next.
    5. The strategic value for the company.
    Rules: max 15 words per bullet. Plain language, no jargon. Use "- " prefix.

    Report:
    {full_report[:6000]}
    """)
    raw = _call_llm(model, prompt)
    bullets = []
    for line in raw.splitlines():
        line = line.strip().lstrip("-•*").strip()
        if line:
            bullets.append(line)
    return bullets[:6]


# ── Report section extractor ────────────────────────────────────────────────

def _extract_sections(report_text: str) -> dict[str, str]:
    """Parse the markdown report into named sections by ## and ### headings."""
    sections: dict[str, str] = {}
    # Split on ## headings (level 2)
    parts = re.split(r'\n(## \d+\.\s+[^\n]+)', report_text)
    for i in range(1, len(parts), 2):
        heading = parts[i].strip().lstrip("#").strip()
        body = parts[i + 1] if i + 1 < len(parts) else ""
        sections[heading] = body.strip()

    # Also extract ### sub-sections within each ## section
    for key, body in list(sections.items()):
        sub_parts = re.split(r'\n(### \d+\.\d+\s+[^\n]+)', body)
        for j in range(1, len(sub_parts), 2):
            sub_heading = sub_parts[j].strip().lstrip("#").strip()
            sub_body = sub_parts[j + 1] if j + 1 < len(sub_parts) else ""
            sections[sub_heading] = sub_body.strip()

    return sections


def _find_section(sections: dict[str, str], *keywords: str) -> str:
    """Find a section whose key contains ALL given keywords (case-insensitive)."""
    for key, body in sections.items():
        if all(kw.lower() in key.lower() for kw in keywords):
            return body
    return ""


# ── Slide builders ──────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _add_bg(slide, color=CMC_LIGHT_BG):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_bottom_bar(slide):
    """Add a thin brand-colour bar at the bottom."""
    left, top = Inches(0), SLIDE_H - Inches(0.35)
    shape = slide.shapes.add_shape(
        1, left, top, SLIDE_W, Inches(0.35)  # MSO_SHAPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CMC_BLUE
    shape.line.fill.background()
    # Footer text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CMC Global  |  AI Pilot Sprint 0 Report  |  April 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = CMC_WHITE
    p.alignment = PP_ALIGN.CENTER


def _add_logo(slide, left=Inches(0.4), top=Inches(0.25), height=Inches(0.55)):
    """Place logo top-left."""
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left, top, height=height)


def _add_title_text(slide, title: str, subtitle: str = "",
                    top=Inches(0.2), left=Inches(2.0)):
    """Add title + optional subtitle on a content slide."""
    txBox = slide.shapes.add_textbox(left, top, Inches(10), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = CMC_DARK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = CMC_GRAY


def _add_bullets(slide, bullets: list[str],
                 left=Inches(0.6), top=Inches(1.3),
                 width=Inches(5.5), height=Inches(5.2),
                 font_size=Pt(14)):
    """Add bullet-point text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {b}"
        p.font.size = font_size
        p.font.color.rgb = CMC_DARK
        p.space_after = Pt(8)
        p.line_spacing = Pt(22)


def _add_chart_image(slide, chart_path: str | Path,
                     left=Inches(6.4), top=Inches(1.2),
                     width=None, height=Inches(5.0)):
    """Place a chart image on the right side."""
    p = Path(chart_path)
    if p.exists():
        slide.shapes.add_picture(str(p), left, top, width=width, height=height)


# ── Slide assembly ──────────────────────────────────────────────────────────

def build_presentation(model: str, output_path: str) -> str:
    report_path = _find_report_md()
    print(f"  Using report: {report_path}")
    report_text = report_path.read_text(encoding="utf-8")
    sec = _extract_sections(report_text)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ================================================================
    # SLIDE 1 — Title
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl, CMC_DARK)

    if LOGO_PATH.exists():
        sl.shapes.add_picture(str(LOGO_PATH), Inches(0.5), Inches(0.4), height=Inches(0.7))

    txBox = sl.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI-Powered Development Pilot"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = CMC_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Sprint 0 Performance Report"
    p2.font.size = Pt(28)
    p2.font.color.rgb = CMC_BLUE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "April 2026  |  CMC Global"
    p3.font.size = Pt(16)
    p3.font.color.rgb = CMC_WHITE
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(24)

    bar = sl.shapes.add_shape(1, Inches(3), Inches(4.8), Inches(7), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CMC_BLUE
    bar.line.fill.background()

    # ================================================================
    # SLIDE 2 — Project Overview & Team (hardcoded per user request)
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "Project Overview")
    _add_bottom_bar(sl)

    overview_bullets = [
        "Apply AI tools in real-life software development to increase efficiency and save cost",
        "Reduce project duration through AI-assisted coding and automation",
        "Document team workflows to transfer knowledge to other projects",
        "Develop internal modules and best practices for the company",
        "Pilot ran 32 AI-assisted sessions across 6 staff members in Sprint 0",
    ]
    _add_bullets(sl, overview_bullets, width=Inches(5.5))

    team_data = [
        ("Frontend Developer", "1 yr exp"),
        ("Backend Developer", "2 yrs exp"),
        ("Tech Lead", "3 yrs exp"),
        ("AI Engineer + Project Ordinator", "2 yrs exp"),
        ("Project Manager", "5 yrs exp"),
        ("Tester", "2 yrs exp"),
        ("QA", "5 yrs exp"),
    ]
    txBox = sl.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Team Composition (7 members)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CMC_BLUE
    p.space_after = Pt(12)
    for role, exp in team_data:
        p = tf.add_paragraph()
        run1 = p.add_run()
        run1.text = f"\u2022  {role}"
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = CMC_DARK
        run2 = p.add_run()
        run2.text = f"  — {exp}"
        run2.font.size = Pt(13)
        run2.font.color.rgb = CMC_GRAY
        p.space_after = Pt(4)

    # ================================================================
    # SLIDE 3 — Executive Summary (KPI chart)
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "Executive Summary", "Sprint 0 Key Results")
    _add_bottom_bar(sl)

    exec_text = _find_section(sec, "Executive Summary")
    exec_bullets = condense(model, exec_text, "Executive Summary", max_bullets=6)
    _add_bullets(sl, exec_bullets, width=Inches(5.5))

    _add_chart_image(sl, CHARTS_DIR / "03_kpi_summary.png",
                     left=Inches(6.4), top=Inches(1.3), height=Inches(4.8))

    # ================================================================
    # SLIDE 4 — Staff Effectiveness + Per-Person Insights
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "Team Performance", "Time Saved per Team Member")
    _add_bottom_bar(sl)

    # Combine per-person insights from section 4.1 for richer LLM input
    person_insights = _find_section(sec, "Per-Person", "Insights")
    staff_text = person_insights or _find_section(sec, "Team Member Subjective")
    staff_bullets = condense(model, staff_text, "Team Performance", max_bullets=6)
    _add_bullets(sl, staff_bullets, width=Inches(5.2))

    _add_chart_image(sl, CHARTS_DIR / "02_staff_ai_effectiveness.png",
                     left=Inches(5.8), top=Inches(1.2), height=Inches(5.2))

    # ================================================================
    # SLIDE 5 — SDLC Tasks by Stage
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "AI Impact Across Development Stages")
    _add_bottom_bar(sl)

    # Pull the closing paragraph of executive summary about strongest adoption
    # plus any SDLC-related content from the report
    sdlc_text = exec_text  # exec summary has SDLC adoption insights
    sdlc_bullets = condense(model, sdlc_text,
                            "AI Impact by Development Stage", max_bullets=6)
    _add_bullets(sl, sdlc_bullets, width=Inches(4.8))

    _add_chart_image(sl, CHARTS_DIR / "01_sdlc_tasks_by_stage.png",
                     left=Inches(5.5), top=Inches(1.2), width=Inches(7.3), height=None)

    # ================================================================
    # SLIDE 6 — AI Tools Comparison
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "AI Tools Performance", "Which Tools Delivered the Most Value")
    _add_bottom_bar(sl)

    tools_text = _find_section(sec, "Tool-Specific", "Findings")
    tools_bullets = condense(model, tools_text, "AI Tools Performance", max_bullets=6)
    _add_bullets(sl, tools_bullets, width=Inches(5.2))

    _add_chart_image(sl, CHARTS_DIR / "04_est_actual_tool.png",
                     left=Inches(5.8), top=Inches(1.2), height=Inches(5.2))

    # ================================================================
    # SLIDE 7 — User Satisfaction / Ratings
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "User Satisfaction", "Rating Distribution by AI Tool")
    _add_bottom_bar(sl)

    # Combine tool findings + exec summary rating info
    rating_text = (
        tools_text + "\n\n"
        + "From Executive Summary: " + exec_text
    )
    rating_bullets = condense(model, rating_text[:3000],
                              "User Satisfaction Ratings", max_bullets=6)
    _add_bullets(sl, rating_bullets, width=Inches(5.2))

    _add_chart_image(sl, CHARTS_DIR / "06_rating_distribution.png",
                     left=Inches(5.8), top=Inches(1.2), height=Inches(5.0))

    # ================================================================
    # SLIDE 8 — Instruction Quality / Top Errors
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "Instruction Quality Analysis",
                    "Most Common Mistakes When Asking AI for Help")
    _add_bottom_bar(sl)

    errors_text = _find_section(sec, "Team-Wide", "Prompt Quality")
    errors_bullets = condense(model, errors_text,
                              "Instruction Quality Analysis", max_bullets=6)
    _add_bullets(sl, errors_bullets, width=Inches(5.2))

    _add_chart_image(sl, CHARTS_DIR / "07_top_errors.png",
                     left=Inches(5.8), top=Inches(1.2), height=Inches(5.0))

    # ================================================================
    # SLIDE 9 — Per-Person Prompt Analysis (from section 2.2)
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "Individual AI Usage Insights",
                    "What Each Team Member Does Well & Where to Improve")
    _add_bottom_bar(sl)

    person_prompt_text = _find_section(sec, "Per-Person", "Prompt")
    person_prompt_bullets = condense(
        model, person_prompt_text[:4000],
        "Individual AI Usage Insights", max_bullets=8
    )

    mid = len(person_prompt_bullets) // 2
    _add_bullets(sl, person_prompt_bullets[:mid],
                 left=Inches(0.6), top=Inches(1.3),
                 width=Inches(5.8), height=Inches(5.0), font_size=Pt(13))
    _add_bullets(sl, person_prompt_bullets[mid:],
                 left=Inches(6.8), top=Inches(1.3),
                 width=Inches(5.8), height=Inches(5.0), font_size=Pt(13))

    # ================================================================
    # SLIDE 10 — Error Heatmap + AI vs Self-Reported
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "Deeper Insights",
                    "Error Patterns by Person & AI vs Self-Assessment")
    _add_bottom_bar(sl)

    _add_chart_image(sl, CHARTS_DIR / "08_error_heatmap.png",
                     left=Inches(0.4), top=Inches(1.3), width=Inches(6.2), height=None)
    _add_chart_image(sl, CHARTS_DIR / "09_user_vs_ai_comparison.png",
                     left=Inches(6.8), top=Inches(1.3), width=Inches(6.0), height=None)

    # ================================================================
    # SLIDE 11 — Team Biases & Adaptation (from section 4.2)
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "Team Adaptation & Biases",
                    "Who Adapts Fastest and Common Self-Assessment Gaps")
    _add_bottom_bar(sl)

    biases_text = _find_section(sec, "Team-Wide", "Patterns")
    biases_bullets = condense(model, biases_text,
                              "Team Adaptation & Biases", max_bullets=7)
    _add_bullets(sl, biases_bullets,
                 left=Inches(0.6), top=Inches(1.3),
                 width=Inches(5.5), height=Inches(5.0))

    _add_chart_image(sl, CHARTS_DIR / "05_est_actual_category.png",
                     left=Inches(6.2), top=Inches(1.2), height=Inches(5.2))

    # ================================================================
    # SLIDE 12 — AI Tool Workarounds & Best Practices (section 3.1 + 3.3)
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "AI Tool Constraints & Workarounds",
                    "Lessons Learned from Real Usage")
    _add_bottom_bar(sl)

    limitations_text = _find_section(sec, "Known Limitations")
    workarounds_text = _find_section(sec, "Recommended Workarounds")
    combined_tools = (limitations_text + "\n\n" + workarounds_text)[:4000]
    tool_wp_bullets = condense(model, combined_tools,
                               "AI Tool Constraints & Practical Workarounds",
                               max_bullets=8)

    mid = len(tool_wp_bullets) // 2
    txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Known Constraints"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = CMC_BLUE

    _add_bullets(sl, tool_wp_bullets[:mid],
                 left=Inches(0.6), top=Inches(1.7),
                 width=Inches(5.5), height=Inches(4.5), font_size=Pt(13))

    txBox2 = sl.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Practical Workarounds"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = CMC_BLUE

    _add_bullets(sl, tool_wp_bullets[mid:],
                 left=Inches(6.8), top=Inches(1.7),
                 width=Inches(5.5), height=Inches(4.5), font_size=Pt(13))

    # ================================================================
    # SLIDE 13 — Recommendations (sections 5.1 + 5.2)
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "Recommendations", "What We Should Do Next")
    _add_bottom_bar(sl)

    rec_text = _find_section(sec, "Recommendations")
    rec_bullets = condense(model, rec_text, "Recommendations", max_bullets=8)

    mid = len(rec_bullets) // 2
    left_bullets = rec_bullets[:mid] if mid else rec_bullets
    right_bullets = rec_bullets[mid:] if mid else []

    txBox = sl.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Immediate Actions"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = CMC_BLUE

    _add_bullets(sl, left_bullets,
                 left=Inches(0.6), top=Inches(1.7),
                 width=Inches(5.5), height=Inches(4.5), font_size=Pt(13))

    txBox2 = sl.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Medium-Term Improvements"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = CMC_BLUE

    _add_bullets(sl, right_bullets,
                 left=Inches(6.8), top=Inches(1.7),
                 width=Inches(5.5), height=Inches(4.5), font_size=Pt(13))

    # ================================================================
    # SLIDE 14 — Training Plan (section 5.3)
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl)
    _add_logo(sl)
    _add_title_text(sl, "Training Plan",
                    "Targeted Skill Development for Each Team Member")
    _add_bottom_bar(sl)

    training_text = _find_section(sec, "Training")
    training_bullets = condense(model, training_text,
                                "Training Plan per Team Member", max_bullets=7)
    _add_bullets(sl, training_bullets,
                 left=Inches(0.6), top=Inches(1.3),
                 width=Inches(12.0), height=Inches(5.0), font_size=Pt(14))

    # ================================================================
    # SLIDE 15 — Conclusion
    # ================================================================
    sl = prs.slides.add_slide(blank_layout)
    _add_bg(sl, CMC_DARK)

    if LOGO_PATH.exists():
        sl.shapes.add_picture(str(LOGO_PATH), Inches(0.5), Inches(0.4), height=Inches(0.6))

    txBox = sl.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusion & Next Steps"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CMC_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Feed the closing note + full report to LLM for a rich conclusion
    closing_note = report_text.split("## Closing note")[-1] if "## Closing note" in report_text else ""
    conclusion_input = closing_note + "\n\n" + report_text[:6000]
    conclusion_bullets = generate_conclusion(model, conclusion_input)

    txBox2 = sl.shapes.add_textbox(Inches(2.0), Inches(2.8), Inches(9.0), Inches(3.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(conclusion_bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"\u2022  {b}"
        p.font.size = Pt(17)
        p.font.color.rgb = CMC_WHITE
        p.space_after = Pt(10)
        p.line_spacing = Pt(28)

    txBox3 = sl.shapes.add_textbox(Inches(1.0), Inches(6.2), Inches(11), Inches(0.6))
    tf3 = txBox3.text_frame
    p = tf3.paragraphs[0]
    p.text = "Thank you  |  CMC Global  |  Aspire to Inspire the Digital World"
    p.font.size = Pt(14)
    p.font.color.rgb = CMC_BLUE
    p.alignment = PP_ALIGN.CENTER

    # ── Save ────────────────────────────────────────────────────────────
    prs.save(output_path)
    return output_path


# ── CLI ─────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description="Generate Sprint 0 AI report slides")
    ap.add_argument("-o", "--output", default="sprint_0_report.pptx",
                    help="Output .pptx path (default: sprint_0_report.pptx)")
    ap.add_argument("--model", default="gpt-5.4-mini",
                    help="LLM model for content condensing (default: gpt-5.4-mini)")
    args = ap.parse_args()

    print(f"Generating slides with model={args.model} ...")
    out = build_presentation(args.model, args.output)
    print(f"Done! Saved to: {out}")


if __name__ == "__main__":
    main()
